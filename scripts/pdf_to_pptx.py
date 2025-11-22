#!/usr/bin/env python3
"""
PDF to PPTX conversion script using PyMuPDF and python-pptx
This script converts a PDF file to PPTX format with better formatting preservation.
"""

import sys
import os
from pathlib import Path

try:
    import fitz  # PyMuPDF
except ImportError:
    print("ERROR: PyMuPDF library not installed. Install it with: pip install pymupdf", file=sys.stderr)
    sys.exit(1)

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
except ImportError:
    print("ERROR: python-pptx library not installed. Install it with: pip install python-pptx", file=sys.stderr)
    sys.exit(1)

def convert_pdf_to_pptx(pdf_path: str, pptx_path: str) -> bool:
    """
    Convert PDF to PPTX using PyMuPDF and python-pptx.
    Each PDF page becomes a PowerPoint slide.
    
    Args:
        pdf_path: Path to input PDF file
        pptx_path: Path to output PPTX file
        
    Returns:
        True if conversion successful, False otherwise
    """
    try:
        # Check if PDF file exists
        if not os.path.exists(pdf_path):
            print(f"ERROR: PDF file not found: {pdf_path}", file=sys.stderr)
            return False
        
        # Open PDF
        pdf_doc = fitz.open(pdf_path)
        
        # Create PowerPoint presentation
        prs = Presentation()
        prs.slide_width = Inches(5)
        prs.slide_height = Inches(7.5)
        
        # Convert each PDF page to a slide
        for page_num in range(len(pdf_doc)):
            page = pdf_doc[page_num]
            
            # Get PDF page dimensions
            page_rect = page.rect
            pdf_width = page_rect.width
            pdf_height = page_rect.height
            pdf_aspect_ratio = pdf_width / pdf_height
            
            # Create a new slide
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
            
            # Calculate slide aspect ratio
            slide_aspect_ratio = prs.slide_width / prs.slide_height
            
            # Track if content was added
            content_added = False
            img_path = None
            
            # Try to convert page to image and add to slide
            try:
                # Render page as image (300 DPI for high quality)
                # Higher DPI = better quality but larger file size
                zoom = 300.0 / 72.0  # 300 DPI for crisp, high-quality images
                mat = fitz.Matrix(zoom, zoom)
                pix = page.get_pixmap(matrix=mat, alpha=False)
                
                # Save image temporarily
                import tempfile
                temp_dir = tempfile.gettempdir()
                img_path = os.path.join(temp_dir, f"pdf_page_{page_num}_{os.getpid()}.png")
                pix.save(img_path)
                
                # Verify image was created and has content
                if os.path.exists(img_path) and os.path.getsize(img_path) > 100:  # At least 100 bytes
                    # Calculate image dimensions - fill width completely, remove left/right margins
                    # Convert PDF dimensions from points to inches (72 points = 1 inch)
                    pdf_width_inches = pdf_width / 72.0
                    pdf_height_inches = pdf_height / 72.0
                    
                    # Fit to full width of slide (no left/right margins)
                    img_width = prs.slide_width
                    # Scale height proportionally to maintain aspect ratio
                    img_height = (pdf_height_inches / pdf_width_inches) * img_width
                    
                    # Center vertically if image is shorter than slide height
                    img_left = 0  # No left margin - starts at left edge
                    if img_height <= prs.slide_height:
                        img_top = (prs.slide_height - img_height) / 2  # Center vertically
                    else:
                        # If image is taller than slide, start at top (may crop bottom)
                        img_top = 0
                        # Optionally scale down to fit height if needed
                        if img_height > prs.slide_height:
                            scale = prs.slide_height / img_height
                            img_height = prs.slide_height
                            img_width = img_width * scale
                            img_left = (prs.slide_width - img_width) / 2  # Re-center horizontally
                    
                    # Add image to slide with proper aspect ratio
                    try:
                        slide.shapes.add_picture(img_path, img_left, img_top, img_width, img_height)
                        content_added = True
                        print(f"Added page {page_num + 1} as image to slide (size: {os.path.getsize(img_path)} bytes, aspect ratio preserved, minimal margins)", file=sys.stderr)
                    except Exception as add_error:
                        print(f"Error adding image to slide: {add_error}", file=sys.stderr)
                        content_added = False
                else:
                    print(f"Image file invalid or too small: {img_path}", file=sys.stderr)
                    content_added = False
                    
            except Exception as img_error:
                print(f"Warning: Could not convert page {page_num + 1} to image: {img_error}", file=sys.stderr)
                content_added = False
            
            # If image didn't work, try text extraction
            if not content_added:
                try:
                    text = page.get_text()
                    if text and text.strip():
                        # Add text box
                        left = Inches(0.5)
                        top = Inches(0.5)
                        width = Inches(9)
                        height = Inches(6.5)
                        text_box = slide.shapes.add_textbox(left, top, width, height)
                        text_frame = text_box.text_frame
                        text_frame.word_wrap = True
                        text_frame.margin_bottom = Inches(0.1)
                        text_frame.margin_top = Inches(0.1)
                        text_frame.margin_left = Inches(0.1)
                        text_frame.margin_right = Inches(0.1)
                        
                        # Add text content (split by lines)
                        lines = [line.strip() for line in text.strip().split('\n') if line.strip()]
                        if lines:
                            for i, line in enumerate(lines[:50]):  # Limit to 50 lines per slide
                                p = text_frame.add_paragraph()
                                p.text = line
                                p.font.size = Pt(14)
                                if i == 0:
                                    p.font.bold = True
                                    p.font.size = Pt(18)
                            content_added = True
                            print(f"Added page {page_num + 1} text to slide ({len(lines)} lines)", file=sys.stderr)
                except Exception as text_error:
                    print(f"Error extracting text from page {page_num + 1}: {text_error}", file=sys.stderr)
            
            # If still no content, add a placeholder
            if not content_added:
                try:
                    left = Inches(0.5)
                    top = Inches(0.5)
                    width = Inches(9)
                    height = Inches(6.5)
                    text_box = slide.shapes.add_textbox(left, top, width, height)
                    text_frame = text_box.text_frame
                    p = text_frame.add_paragraph()
                    p.text = f"Page {page_num + 1} of {len(pdf_doc)}"
                    p.font.size = Pt(24)
                    p.font.bold = True
                    print(f"Added placeholder for page {page_num + 1}", file=sys.stderr)
                except Exception as placeholder_error:
                    print(f"Error adding placeholder: {placeholder_error}", file=sys.stderr)
            
            # Clean up temp image if it exists
            if img_path and os.path.exists(img_path):
                try:
                    os.unlink(img_path)
                except:
                    pass
        
        # Close PDF
        pdf_doc.close()
        
        # Save PowerPoint presentation
        prs.save(pptx_path)
        
        # Check if output file was created
        if os.path.exists(pptx_path):
            print(f"SUCCESS: Converted {pdf_path} to {pptx_path}")
            return True
        else:
            print(f"ERROR: Output file was not created: {pptx_path}", file=sys.stderr)
            return False
            
    except Exception as e:
        print(f"ERROR: Conversion failed: {str(e)}", file=sys.stderr)
        import traceback
        traceback.print_exc()
        return False

if __name__ == "__main__":
    if len(sys.argv) != 3:
        print("Usage: python pdf_to_pptx.py <input_pdf> <output_pptx>", file=sys.stderr)
        sys.exit(1)
    
    input_pdf = sys.argv[1]
    output_pptx = sys.argv[2]
    
    success = convert_pdf_to_pptx(input_pdf, output_pptx)
    sys.exit(0 if success else 1)

